from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib.colors import HexColor, white
from reportlab.lib.pagesizes import landscape
from reportlab.lib.utils import ImageReader
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.pdfgen import canvas


ROOT = Path(__file__).resolve().parent
OUTPUT_DIR = ROOT / "outputs"
PLOTS_DIR = OUTPUT_DIR / "plots"
DELIVERABLES_DIR = ROOT / "deliverables"

PPTX_PATH = DELIVERABLES_DIR / "athena_ot_poc_report.pptx"
PDF_PATH = DELIVERABLES_DIR / "athena_ot_poc_report.pdf"

SLIDE_W = 13.333
SLIDE_H = 7.5

COLORS = {
    "navy": RGBColor(12, 25, 59),
    "blue": RGBColor(47, 107, 255),
    "sky": RGBColor(226, 237, 255),
    "text": RGBColor(35, 42, 52),
    "muted": RGBColor(96, 109, 128),
    "line": RGBColor(213, 222, 235),
    "white": RGBColor(255, 255, 255),
}


def ensure_dirs() -> None:
    DELIVERABLES_DIR.mkdir(exist_ok=True)


def load_metrics() -> pd.DataFrame:
    return pd.read_csv(OUTPUT_DIR / "metrics.csv")


def load_eda() -> dict:
    return pd.read_json(OUTPUT_DIR / "eda_summary.json", typ="series").to_dict()


def load_correlations() -> pd.DataFrame:
    corr = pd.read_csv(OUTPUT_DIR / "target_correlations.csv")
    corr.columns = ["feature", "correlation"]
    return corr


def model_label(model_name: str) -> str:
    labels = {
        "ridge": "Ridge",
        "linear_regression": "Linear Regression",
        "extra_trees": "Extra Trees",
        "hist_gradient_boosting": "HistGradientBoosting",
        "random_forest": "Random Forest",
        "last_value_baseline": "直前値ベースライン",
        "mean_baseline": "平均値ベースライン",
    }
    return labels.get(model_name, model_name)


def build_story() -> list[dict]:
    metrics = load_metrics()
    eda = load_eda()
    corr = load_correlations()

    test_metrics = metrics[metrics["split"] == "test"].sort_values("mae").reset_index(drop=True)
    best = test_metrics.iloc[0]
    naive = test_metrics[test_metrics["model"] == "last_value_baseline"].iloc[0]
    improvement = (naive["mae"] - best["mae"]) / naive["mae"] * 100
    top_corr = corr[corr["feature"] != "OT"].head(3)
    corr_text = " / ".join(
        f"{row.feature} ({row.correlation:.2f})" for row in top_corr.itertuples(index=False)
    )

    result_lines = [
        f"{model_label(row.model)}: MAE {row.mae:.3f} / RMSE {row.rmse:.3f} / R2 {row.r2:.3f}"
        for row in test_metrics.head(5).itertuples(index=False)
    ]

    return [
        {
            "title": "OTの1時間先予測を試したPoC",
            "subtitle": "初学者向けに整理した ETTh1 データの予測課題サマリー",
        },
        {
            "title": "今回やったことを一言でいうと",
            "bullets": [
                "過去のデータを使って、1時間後の OT を予測できるか試した",
                "まずデータの中身を確認し、そのあと予測に使う材料を作った",
                "複数のモデルを比べて、どの方法が一番当たるかを確認した",
            ],
            "accent": (
                f"最良モデルは {model_label(best['model'])}。"
                f" テストでは MAE {best['mae']:.3f}、直前値ベースラインより {improvement:.1f}% 改善"
            ),
        },
        {
            "title": "課題の進め方",
            "bullets": [
                "1. データを読む",
                "2. EDAでデータの状態を確認する",
                "3. 予測に使う特徴量を作る",
                "4. train / valid / test に分ける",
                "5. モデルを学習して性能を比べる",
                "6. 結果を表とグラフでまとめる",
            ],
        },
        {
            "title": "EDAで確認したこと",
            "bullets": [
                f"データ期間は {str(eda['date_min'])[:10]} から {str(eda['date_max'])[:10]}、件数は {int(eda['rows']):,} 行",
                "欠損値は全列で 0 件だったので、まずはそのまま使える状態だった",
                f"目的変数 OT の平均は {eda['target_mean']:.2f}、標準偏差は {eda['target_std']:.2f}",
                f"OT と相関が比較的高かった列は {corr_text}",
            ],
            "image": PLOTS_DIR / "ot_timeseries.png",
        },
        {
            "title": "EDAとは何か",
            "bullets": [
                "EDA は Exploratory Data Analysis の略で、探索的データ分析のこと",
                "いきなりモデルを作らず、まずデータの特徴やクセを理解する作業",
                "欠損、分布、時系列の動き、効きそうな列を把握してから次に進む",
            ],
            "accent": "今回の EDA は、OT が予測できそうかを見極めるための下調べ",
        },
        {
            "title": "予測に使う材料を作った",
            "bullets": [
                "時刻から hour、dayofweek、month、is_weekend を作成",
                "OT の 1時間前、2時間前、24時間前などの lag を作成",
                "直近数時間の平均やばらつきも追加して、流れを捉えやすくした",
                "こうした予測に使う材料を特徴量と呼ぶ",
            ],
            "image": PLOTS_DIR / "feature_importance.png",
        },
        {
            "title": "データの分け方",
            "bullets": [
                "時系列データなので、順番を崩さずに古い順から分割した",
                "学習用 train が 70%、調整用 valid が 15%、最終確認用 test が 15%",
                "未来の情報を見ないようにするのが重要",
            ],
            "accent": "今回の設定は「時刻 T までの情報で T+1 の OT を予測する」",
        },
        {
            "title": "試したモデル",
            "bullets": [
                "比較用として平均値ベースラインと直前値ベースラインを用意した",
                "学習モデルとして Linear Regression、Ridge、Random Forest などを試した",
                "単純な方法より良いかどうかを確認するのが PoC では大事",
            ],
        },
        {
            "title": "結果",
            "bullets": result_lines,
            "accent": "今回は Ridge が最も安定して高精度だった",
        },
        {
            "title": "この結果から言えること",
            "bullets": [
                "OT の 1時間先予測は十分現実的で、PoC として成立している",
                "複雑すぎるモデルでなくても、特徴量設計でかなり良い精度が出た",
                "まずはシンプルな手法でベースを作る進め方は妥当だった",
            ],
            "image": PLOTS_DIR / "test_predictions.png",
        },
        {
            "title": "次にやるとよいこと",
            "bullets": [
                "6時間先や24時間先など、より難しい予測にも広げる",
                "LightGBM などの別モデルも試して比較を増やす",
                "複数期間のバックテストで、時期による性能差も確認する",
            ],
            "accent": "今回の成果は「1時間先予測の最初の土台づくり」として使える",
        },
    ]


def style_text_frame(text_frame, font_size: int, color: RGBColor, bold: bool = False) -> None:
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.name = "Yu Gothic"


def add_title(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(11.5), Inches(0.7))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    p.alignment = PP_ALIGN.LEFT
    style_text_frame(tf, 26, COLORS["navy"], bold=True)


def add_bullets(slide, bullets: list[str], left: float = 0.85, top: float = 1.45, width: float = 5.0) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(4.8))
    tf = box.text_frame
    tf.word_wrap = True
    for index, bullet in enumerate(bullets):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.space_after = Pt(8)
    style_text_frame(tf, 18, COLORS["text"])


def add_accent(slide, text: str) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.85),
        Inches(5.95),
        Inches(11.6),
        Inches(0.7),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["sky"]
    shape.line.color.rgb = COLORS["line"]
    tf = shape.text_frame
    tf.text = text
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    style_text_frame(tf, 16, COLORS["navy"], bold=True)


def add_image(slide, path: Path, left: float = 6.1, top: float = 1.45, width: float = 6.5) -> None:
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))


def create_pptx(story: list[dict]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    for index, section in enumerate(story):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = COLORS["white"]

        top_band = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0),
            Inches(0),
            Inches(SLIDE_W),
            Inches(0.22),
        )
        top_band.fill.solid()
        top_band.fill.fore_color.rgb = COLORS["blue"]
        top_band.line.fill.background()

        add_title(slide, section["title"])

        if index == 0:
            hero = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                Inches(0.85),
                Inches(1.65),
                Inches(11.6),
                Inches(3.1),
            )
            hero.fill.solid()
            hero.fill.fore_color.rgb = COLORS["navy"]
            hero.line.fill.background()
            tf = hero.text_frame
            tf.text = section["title"]
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            style_text_frame(tf, 28, COLORS["white"], bold=True)

            subtitle = slide.shapes.add_textbox(Inches(1.3), Inches(5.1), Inches(10.8), Inches(0.8))
            subtitle_tf = subtitle.text_frame
            subtitle_tf.text = section["subtitle"]
            subtitle_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            style_text_frame(subtitle_tf, 18, COLORS["muted"])
            continue

        if "bullets" in section:
            text_width = 5.0 if "image" in section else 11.0
            add_bullets(slide, section["bullets"], width=text_width)

        if "image" in section:
            add_image(slide, section["image"])

        if "accent" in section:
            add_accent(slide, section["accent"])

    prs.save(PPTX_PATH)


def wrap_text(text: str, width: int = 36) -> list[str]:
    lines: list[str] = []
    current = ""
    for chunk in text.split(" "):
        candidate = chunk if not current else f"{current} {chunk}"
        if len(candidate) <= width:
            current = candidate
        else:
            if current:
                lines.append(current)
            current = chunk
    if current:
        lines.append(current)
    return lines


def create_pdf(story: list[dict]) -> None:
    pdfmetrics.registerFont(UnicodeCIDFont("HeiseiKakuGo-W5"))
    c = canvas.Canvas(str(PDF_PATH), pagesize=landscape((960, 540)))
    width, height = landscape((960, 540))

    for index, section in enumerate(story, start=1):
        c.setFillColor(HexColor("#FFFFFF"))
        c.rect(0, 0, width, height, stroke=0, fill=1)
        c.setFillColor(HexColor("#2F6BFF"))
        c.rect(0, height - 16, width, 16, stroke=0, fill=1)

        c.setFont("HeiseiKakuGo-W5", 24)
        c.setFillColor(HexColor("#0C193B"))
        c.drawString(40, height - 54, section["title"])

        if index == 1:
            c.setFillColor(HexColor("#0C193B"))
            c.roundRect(70, 180, width - 140, 170, 18, stroke=0, fill=1)
            c.setFont("HeiseiKakuGo-W5", 26)
            c.setFillColor(white)
            c.drawCentredString(width / 2, 285, section["title"])
            c.setFont("HeiseiKakuGo-W5", 16)
            c.setFillColor(HexColor("#606D80"))
            c.drawCentredString(width / 2, 145, section["subtitle"])
            c.showPage()
            continue

        y = height - 95
        if "bullets" in section:
            c.setFont("HeiseiKakuGo-W5", 14)
            c.setFillColor(HexColor("#232A34"))
            bullet_x = 52
            for bullet in section["bullets"]:
                c.drawString(bullet_x, y, "-")
                text_x = bullet_x + 16
                wrapped = wrap_text(bullet, width=32 if "image" in section else 70)
                for line in wrapped:
                    c.drawString(text_x, y, line)
                    y -= 20
                y -= 8

        if "image" in section and section["image"].exists():
            image = ImageReader(str(section["image"]))
            c.drawImage(image, 500, 120, width=400, height=250, preserveAspectRatio=True, mask="auto")

        if "accent" in section:
            c.setFillColor(HexColor("#E2EDFF"))
            c.roundRect(50, 36, width - 100, 48, 10, stroke=0, fill=1)
            c.setFillColor(HexColor("#0C193B"))
            c.setFont("HeiseiKakuGo-W5", 13)
            c.drawCentredString(width / 2, 55, section["accent"])

        c.showPage()

    c.save()


def main() -> None:
    ensure_dirs()
    story = build_story()
    create_pptx(story)
    create_pdf(story)
    print(f"Created: {PPTX_PATH}")
    print(f"Created: {PDF_PATH}")


if __name__ == "__main__":
    main()
